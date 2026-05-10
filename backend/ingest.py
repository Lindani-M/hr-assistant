import io
import os
import json
import logging
import requests
from msal import ConfidentialClientApplication
from bs4 import BeautifulSoup
from openai import AzureOpenAI
from dotenv import load_dotenv
import pdfplumber
from pptx import Presentation
from docx import Document

load_dotenv()

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s - %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger(__name__)

# ── Credentials ───────────────────────────────────────────────────────────────
SP_TENANT_ID     = os.getenv("SP_TENANT_ID")
SP_CLIENT_ID     = os.getenv("SP_CLIENT_ID")
SP_CLIENT_SECRET = os.getenv("SP_CLIENT_SECRET")
SP_HOSTNAME      = os.getenv("SP_HOSTNAME")
SP_SITE_PATH     = os.getenv("SP_SITE_PATH")

AZURE_OAI_ENDPOINT = os.getenv("AZURE_OAI_ENDPOINT")
AZURE_OAI_KEY      = os.getenv("AZURE_OAI_KEY")
EMBED_DEPLOYMENT   = os.getenv("EMBED_DEPLOYMENT")

SEARCH_ENDPOINT    = os.getenv("SEARCH_ENDPOINT")
SEARCH_KEY         = os.getenv("SEARCH_KEY")
SEARCH_INDEX       = os.getenv("SEARCH_INDEX")
SEARCH_API_VERSION = os.getenv("SEARCH_API_VERSION")

# Folder path (relative to drive root) to scan for Word/PowerPoint/PDF files
SP_DOCS_FOLDER = os.getenv("SP_DOCS_FOLDER", "Documents")

# ── Azure OpenAI client ───────────────────────────────────────────────────────
oai_client = AzureOpenAI(
    azure_endpoint=AZURE_OAI_ENDPOINT,
    api_key=AZURE_OAI_KEY,
    api_version="2024-02-01",
)


# ══════════════════════════════════════════════════════════════════════════════
# STEP 1 — SCRAPE SHAREPOINT
# ══════════════════════════════════════════════════════════════════════════════

def get_sp_token():
    app = ConfidentialClientApplication(
        client_id=SP_CLIENT_ID,
        client_credential=SP_CLIENT_SECRET,
        authority=f"https://login.microsoftonline.com/{SP_TENANT_ID}",
    )
    result = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    if "access_token" not in result:
        raise Exception(f"SP Auth failed: {result.get('error_description')}")
    return result["access_token"]


def get_site_id(headers: dict) -> str:
    """Resolve the Graph site-id from SP_HOSTNAME and SP_SITE_PATH."""
    site_resp = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{SP_HOSTNAME}:{SP_SITE_PATH}",
        headers=headers,
    ).json()
    if "id" not in site_resp:
        raise Exception(f"Could not resolve site id: {site_resp.get('error', site_resp)}")
    return site_resp["id"]


def scrape_all_pages(sp_token, site_id):
    headers = {"Authorization": f"Bearer {sp_token}"}

    logger.info("Scraping SharePoint site — path=%s, site_id=%s", SP_SITE_PATH, site_id)

    pages = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{site_id}/pages",
        headers=headers,
    ).json().get("value", [])

    logger.info("Found %d pages", len(pages))

    documents = []

    for page in pages:
        page_id   = page["id"]
        page_name = page.get("name", "")
        page_url  = page.get("webUrl", "")

        page_data = requests.get(
            f"https://graph.microsoft.com/beta/sites/{site_id}/pages/{page_id}"
            f"/microsoft.graph.sitePage?$expand=canvasLayout",
            headers=headers,
        ).json()

        title = page_data.get("title", page_name)
        desc  = page_data.get("description", "").strip()

        wp_resp = requests.get(
            f"https://graph.microsoft.com/beta/sites/{site_id}/pages/{page_id}"
            f"/microsoft.graph.sitePage/webParts",
            headers=headers,
        )

        text_blocks = []
        if desc:
            text_blocks.append(desc)

        if wp_resp.status_code == 200:
            for wp in wp_resp.json().get("value", []):
                html = wp.get("innerHtml", "") or wp.get("innerHTML", "")
                if html:
                    soup = BeautifulSoup(html, "html.parser")
                    text = soup.get_text(separator=" ", strip=True)
                    if text and len(text) > 10:
                        text_blocks.append(text)
                else:
                    props = (wp.get("data") or {}).get("properties") or {}
                    for key in ["title", "text", "content", "description"]:
                        val = props.get(key, "")
                        if val and isinstance(val, str) and len(val) > 10:
                            text_blocks.append(val)
                            break

        full_text = "\n\n".join(text_blocks)
        if full_text.strip():
            documents.append({
                "page_id":     page_id,
                "title":       title,
                "url":         page_url,
                "full_text":   full_text,
                "source_type": "page",
            })
            logger.info(
                "✅ Page '%s' — %d chars, %d blocks | %s",
                title, len(full_text), len(text_blocks), page_url,
            )

    logger.info("Total pages scraped: %d", len(documents))
    return documents


# ══════════════════════════════════════════════════════════════════════════════
# STEP 1b — SCRAPE SHAREPOINT DRIVE FILES (PDF / DOCX / PPTX)
# ══════════════════════════════════════════════════════════════════════════════

SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".pptx")


def _get_drive_id(headers: dict, site_id: str) -> str:
    resp = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive",
        headers=headers,
    ).json()
    return resp["id"]


def _list_files_recursive(headers: dict, drive_id: str, folder_path: str) -> list:
    """Return all drive items whose extension is in SUPPORTED_EXTENSIONS."""
    found = []
    url = (
        f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
        f"/root:/{folder_path}:/children?$top=999"
    )
    while url:
        resp = requests.get(url, headers=headers)
        if resp.status_code != 200:
            logger.warning("Drive list failed for '%s': HTTP %d", folder_path, resp.status_code)
            break
        data = resp.json()
        for item in data.get("value", []):
            if "folder" in item:
                sub = f"{folder_path}/{item['name']}"
                found.extend(_list_files_recursive(headers, drive_id, sub))
            elif "file" in item:
                if item["name"].lower().endswith(SUPPORTED_EXTENSIONS):
                    found.append(item)
        url = data.get("@odata.nextLink")
    return found


def _list_files_from_root(headers: dict, drive_id: str) -> list:
    """Scan drive root + all subfolders recursively."""
    found = []
    url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children?$top=999"
    while url:
        resp = requests.get(url, headers=headers)
        if resp.status_code != 200:
            logger.warning("Root list failed: HTTP %d", resp.status_code)
            break
        data = resp.json()
        for item in data.get("value", []):
            if "folder" in item:
                logger.info("  📁 Scanning subfolder: '%s'", item["name"])
                found.extend(_list_files_recursive(headers, drive_id, item["name"]))
            elif "file" in item:
                if item["name"].lower().endswith(SUPPORTED_EXTENSIONS):
                    logger.info("  📄 Found file at root: '%s'", item["name"])
                    found.append(item)
        url = data.get("@odata.nextLink")
    return found


def _download_drive_item(headers: dict, item: dict) -> bytes:
    """Download file bytes via the pre-authenticated download URL when available."""
    download_url = item.get("@microsoft.graph.downloadUrl")
    if download_url:
        resp = requests.get(download_url, timeout=120)
    else:
        drive_id = item["parentReference"]["driveId"]
        item_id  = item["id"]
        resp = requests.get(
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{item_id}/content",
            headers=headers,
            allow_redirects=True,
            timeout=120,
        )
    resp.raise_for_status()
    return resp.content


def _extract_pdf(file_bytes: bytes) -> str:
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text = page.extract_text(
                x_tolerance_ratio=0.15,
                y_tolerance=3,
                use_text_flow=True,
            ) or ""
            if not text.strip():
                text = page.extract_text(
                    x_tolerance=1.5,
                    y_tolerance=3,
                ) or ""
            parts.append(text.strip())
    return "\n\n".join(parts)


def _extract_docx(file_bytes: bytes) -> str:
    doc = Document(io.BytesIO(file_bytes))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = []
    for idx, slide in enumerate(prs.slides, 1):
        shapes_text = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if shapes_text:
            slide_texts.append(f"[Slide {idx}]\n" + "\n".join(shapes_text))
    return "\n\n".join(slide_texts)


def _extract_text(file_bytes: bytes, filename: str) -> str:
    name = filename.lower()
    if name.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    if name.endswith(".docx"):
        return _extract_docx(file_bytes)
    if name.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    return ""


def scrape_drive_documents(sp_token, site_id):
    """Pull text from all supported files on the SharePoint drive (root + all subfolders)."""
    headers  = {"Authorization": f"Bearer {sp_token}"}
    drive_id = _get_drive_id(headers, site_id)
    logger.info("Scanning drive root — drive_id=%s", drive_id)

    items = _list_files_from_root(headers, drive_id)
    logger.info("Found %d supported files total", len(items))

    documents = []
    for item in items:
        name    = item["name"]
        web_url = item.get("webUrl", "")
        logger.info("Processing: %s", name)
        try:
            file_bytes = _download_drive_item(headers, item)
            text       = _extract_text(file_bytes, name)
            if text.strip():
                documents.append({
                    "page_id":     item["id"],
                    "title":       name,
                    "url":         web_url,
                    "full_text":   text,
                    "source_type": "file",
                })
                logger.info("✅ Extracted %d chars from '%s'", len(text), name)
            else:
                logger.warning("⚠️  No text from '%s' — skipping", name)
        except Exception as exc:
            logger.error("❌ Failed '%s': %s", name, exc)

    logger.info("Drive documents extracted: %d", len(documents))
    return documents


# ══════════════════════════════════════════════════════════════════════════════
# STEP 2 — CHUNK
# ══════════════════════════════════════════════════════════════════════════════

def chunk_documents(documents, chunk_size=400, overlap=80):
    chunks = []
    logger.info(
        "Chunking %d documents — chunk_size=%d words, overlap=%d words",
        len(documents), chunk_size, overlap,
    )

    for doc in documents:
        words       = doc["full_text"].split()
        start       = 0
        chunk_index = 0
        logger.info("Chunking '%s' (%d words)", doc["title"], len(words))

        while start < len(words):
            end        = min(start + chunk_size, len(words))
            chunk_text = " ".join(words[start:end])
            chunk_id   = f"{doc['page_id']}_chunk{chunk_index}"

            chunks.append({
                "id":          chunk_id,
                "page_id":     doc["page_id"],
                "title":       doc["title"],
                "url":         doc["url"],
                "source_type": doc.get("source_type", "unknown"),
                "chunk_index": chunk_index,
                "text":        chunk_text,
            })

            logger.debug(
                "  chunk %d: %d words (words %d–%d of %d) id=%s",
                chunk_index, len(chunk_text.split()), start, end, len(words), chunk_id,
            )

            chunk_index += 1
            if end == len(words):
                break
            start += chunk_size - overlap

    avg = sum(len(c["text"].split()) for c in chunks) // max(len(chunks), 1)
    logger.info(
        "Chunking complete: %d chunks across %d pages, avg %d words/chunk",
        len(chunks), len(documents), avg,
    )
    return chunks


# ══════════════════════════════════════════════════════════════════════════════
# STEP 3 — EMBED
# ══════════════════════════════════════════════════════════════════════════════

def embed_chunks(chunks, batch_size=16):
    logger.info(
        "Embedding %d chunks — model=%s, batch_size=%d",
        len(chunks), EMBED_DEPLOYMENT, batch_size,
    )

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        texts = [c["text"] for c in batch]
        resp  = oai_client.embeddings.create(model=EMBED_DEPLOYMENT, input=texts)

        for j, item in enumerate(resp.data):
            batch[j]["embedding"] = item.embedding
            vec = item.embedding
            logger.debug(
                "  chunk '%s' — dim=%d, magnitude=%.4f",
                batch[j]["id"], len(vec), sum(v ** 2 for v in vec) ** 0.5,
            )

        batch_num   = i // batch_size + 1
        total_batch = (len(chunks) - 1) // batch_size + 1
        logger.info("Embedded batch %d/%d", batch_num, total_batch)

    logger.info("All %d chunks embedded", len(chunks))
    return chunks


# ══════════════════════════════════════════════════════════════════════════════
# STEP 4 — CREATE AZURE AI SEARCH INDEX
# ══════════════════════════════════════════════════════════════════════════════

def create_index():
    logger.info("Creating index '%s' at %s", SEARCH_INDEX, SEARCH_ENDPOINT)

    url     = f"{SEARCH_ENDPOINT}/indexes/{SEARCH_INDEX}?api-version={SEARCH_API_VERSION}"
    headers = {"Content-Type": "application/json", "api-key": SEARCH_KEY}

    index_def = {
        "name": SEARCH_INDEX,
        "fields": [
            {"name": "id",          "type": "Edm.String", "key": True,  "filterable": True, "retrievable": True},
            {"name": "page_id",     "type": "Edm.String", "filterable": True, "retrievable": True},
            {"name": "title",       "type": "Edm.String", "searchable": True, "filterable": True, "retrievable": True},
            {"name": "url",         "type": "Edm.String", "filterable": True, "retrievable": True},
            {"name": "source_type", "type": "Edm.String", "filterable": True, "retrievable": True},
            {"name": "chunk_index", "type": "Edm.Int32",  "filterable": True, "sortable": True,  "retrievable": True},
            {"name": "text",        "type": "Edm.String", "searchable": True, "retrievable": True},
            {
                "name": "embedding",
                "type": "Collection(Edm.Single)",
                "searchable": True,
                "retrievable": True,
                "dimensions": 1536,
                "vectorSearchProfile": "myHnswProfile",
            },
        ],
        "vectorSearch": {
            "algorithms": [{"name": "myHnsw", "kind": "hnsw"}],
            "profiles":   [{"name": "myHnswProfile", "algorithm": "myHnsw"}],
        },
    }

    resp = requests.put(url, headers=headers, json=index_def)
    if resp.status_code in (200, 201, 204):
        logger.info("✅ Index '%s' ready (HTTP %d)", SEARCH_INDEX, resp.status_code)
    else:
        logger.error("❌ Index creation failed: %d %s", resp.status_code, resp.text[:300])


# ══════════════════════════════════════════════════════════════════════════════
# STEP 5 — UPLOAD TO AZURE AI SEARCH
# ══════════════════════════════════════════════════════════════════════════════

def upload_chunks(chunks, batch_size=100):
    logger.info("Uploading %d chunks to index '%s'", len(chunks), SEARCH_INDEX)

    url = (
        f"{SEARCH_ENDPOINT}/indexes/{SEARCH_INDEX}"
        f"/docs/index?api-version={SEARCH_API_VERSION}"
    )
    headers = {"Content-Type": "application/json", "api-key": SEARCH_KEY}

    for i in range(0, len(chunks), batch_size):
        batch = chunks[i:i + batch_size]
        body  = {
            "value": [
                {
                    "@search.action": "mergeOrUpload",
                    "id":          c["id"],
                    "page_id":     c["page_id"],
                    "title":       c["title"],
                    "url":         c["url"],
                    "source_type": c.get("source_type", "unknown"),
                    "chunk_index": c["chunk_index"],
                    "text":        c["text"],
                    "embedding":   c["embedding"],
                }
                for c in batch
            ]
        }
        resp      = requests.post(url, headers=headers, json=body)
        batch_num = i // batch_size + 1

        if resp.status_code in (200, 207):
            results = resp.json().get("value", [])
            success = sum(1 for r in results if r.get("status"))
            failed  = len(results) - success
            logger.info("Batch %d: %d uploaded, %d failed", batch_num, success, failed)
            for c in batch:
                logger.debug("  uploaded: %s | %s | chunk %d", c["id"], c["title"], c["chunk_index"])
        else:
            logger.error("Batch %d error: HTTP %d — %s", batch_num, resp.status_code, resp.text[:300])

    logger.info("Upload complete — %d chunks in index '%s'", len(chunks), SEARCH_INDEX)


# ══════════════════════════════════════════════════════════════════════════════
# VERIFY UPLOAD
# ══════════════════════════════════════════════════════════════════════════════

def verify_upload():
    resp = requests.get(
        f"{SEARCH_ENDPOINT}/indexes/{SEARCH_INDEX}/docs"
        f"?api-version={SEARCH_API_VERSION}&$top=1&$select=id,title,chunk_index,embedding",
        headers={"api-key": SEARCH_KEY},
    )
    data = resp.json()

    if "value" in data and data["value"]:
        doc       = data["value"][0]
        embedding = doc.get("embedding", [])
        logger.info(
            "Verification — id=%s, title='%s', chunk=%s, embedding=%s (%d dims)",
            doc["id"], doc["title"], doc["chunk_index"],
            "present" if embedding else "MISSING", len(embedding),
        )
    else:
        logger.warning("Verification: no documents found or error: %s", json.dumps(data))


# ══════════════════════════════════════════════════════════════════════════════
# ENTRYPOINT
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    logger.info("=== SharePoint Ingest Pipeline Starting ===")
    sp_token = get_sp_token()
    _headers = {"Authorization": f"Bearer {sp_token}"}
    site_id  = get_site_id(_headers)

    page_docs  = scrape_all_pages(sp_token, site_id)
    drive_docs = scrape_drive_documents(sp_token, site_id)
    documents  = page_docs + drive_docs
    logger.info(
        "Total source documents: %d (%d pages + %d drive files)",
        len(documents), len(page_docs), len(drive_docs),
    )

    chunks          = chunk_documents(documents)
    embedded_chunks = embed_chunks(chunks)
    create_index()
    upload_chunks(embedded_chunks)
    verify_upload()
    logger.info("=== Ingest Pipeline Complete ===")
